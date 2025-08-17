import json
import base64
import io
import pymupdf
import fitz
import os
from pptx import Presentation
from PIL import Image
from docx import Document
from mistralai import Mistral
import pdfplumber

client = Mistral(api_key=os.getenv("MISTRAL_API_KEY"))

def encode_pdf(pdf_bytes: bytes):
    """Encode PDF bytes to a base64 string."""
    try:
        return base64.b64encode(pdf_bytes).decode("utf-8")
    except Exception as e:
        print(f"Error encoding PDF to base64: {e}")
        return None

def convert_to_pdf(file_bytes: bytes, filename: str):
    extension = filename.lower().split('.')[-1]

    if extension == "pdf":
        return file_bytes

    buffer = io.BytesIO()
    pdf = fitz.open()

    if extension in {"jpg", "jpeg", "png", "gif", "webp", "bmp"}:
        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
        img.save(buffer, format="PDF")
        return buffer.getvalue()

    elif extension in {"txt", "md"}:
        text = file_bytes.decode("utf-8", errors="ignore")
        page = pdf.new_page()
        page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    elif extension in {"doc", "docx"}:
        doc = Document(io.BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs])
        page = pdf.new_page()
        page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    elif extension == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            page = pdf.new_page()
            page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    else:
        raise ValueError(f"Unsupported file type: {extension}")
    
def process_page(idx, ocr_response=None):
    try:
        if ocr_response and hasattr(ocr_response, 'pages') and idx < len(ocr_response.pages):
            return ocr_response.pages[idx].markdown
        else:
            return f"Error: Page {idx + 1} not available in OCR response"
    except Exception as e:
        return f"Error processing page {idx + 1}: {e}"

def extract_text_from_pdf(pdf_bytes: bytes):
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            total_pages = len(pdf.pages)
    except Exception:
        with pymupdf.open(stream=pdf_bytes, filetype="pdf") as doc:
            total_pages = len(doc)

    pdf_bytes = encode_pdf(pdf_bytes)

    try:
        response = client.ocr.process(
            model="mistral-ocr-latest",
            document={
                "type": "document_url",
                "document_url": f"data:application/pdf;base64,{pdf_bytes}"
            },
            include_image_base64=True
        )
    except Exception as e:
        return [f"Error during OCR processing: {e}"]

    page_numbers = list(range(total_pages))

    extracted_text = []
    for idx in page_numbers:
        extracted_text.append(process_page(idx, ocr_response=response))

    return extracted_text

def create_chunks(directory_path: str):
    file_paths = [
        os.path.abspath(os.path.join(directory_path, f))
        for f in os.listdir(directory_path)
        if os.path.isfile(os.path.join(directory_path, f))
    ]
    Chunks = []
    for idx, file_path in enumerate(file_paths):
        filename = os.path.basename(file_path)
        with open(file_path, "rb") as f:
            file_bytes = f.read()
        converted_pdf_bytes = convert_to_pdf(file_bytes, filename)
        print(f"Processing file: {filename}")
        pages = extract_text_from_pdf(converted_pdf_bytes)
        for page in pages:
            Chunks.append({"filename": filename, "page_number": idx + 1,"page_content": page})
    return Chunks