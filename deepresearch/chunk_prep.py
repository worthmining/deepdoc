import json
import base64
import io
import pymupdf
import fitz
import os
from pptx import Presentation
from PIL import Image
from docx import Document
from mistralai import Mistral
import pdfplumber

client = Mistral(api_key=os.getenv("MISTRAL_API_KEY"))

def encode_pdf(pdf_bytes: bytes):
    """Encode PDF bytes to a base64 string."""
    try:
        return base64.b64encode(pdf_bytes).decode("utf-8")
    except Exception as e:
        print(f"Error encoding PDF to base64: {e}")
        return None

def convert_to_pdf(file_bytes: bytes, filename: str):
    extension = filename.lower().split('.')[-1]

    if extension == "pdf":
        return file_bytes

    buffer = io.BytesIO()
    pdf = fitz.open()

    if extension in {"jpg", "jpeg", "png", "gif", "webp", "bmp"}:
        img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
        img.save(buffer, format="PDF")
        return buffer.getvalue()

    elif extension in {"txt", "md"}:
        text = file_bytes.decode("utf-8", errors="ignore")
        lines = text.splitlines()
        pdf = fitz.open()
        max_lines_per_page = 40  # You can adjust this limit

        for i in range(0, len(lines), max_lines_per_page):
            page = pdf.new_page()
            chunk_text = "\n".join(lines[i:i + max_lines_per_page])
            page.insert_text((72, 72), chunk_text)

        pdf.save(buffer)
        return buffer.getvalue()


    elif extension in {"doc", "docx"}:
        doc = Document(io.BytesIO(file_bytes))
        pdf = fitz.open()
        paragraphs = [para.text for para in doc.paragraphs]
        max_paras_per_page = 20  # Adjustable limit

        for i in range(0, len(paragraphs), max_paras_per_page):
            page = pdf.new_page()
            chunk_text = "\n".join(paragraphs[i:i + max_paras_per_page])
            page.insert_text((72, 72), chunk_text)

        pdf.save(buffer)
        return buffer.getvalue()


    elif extension == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            page = pdf.new_page()
            page.insert_text((72, 72), text)
        pdf.save(buffer)
        return buffer.getvalue()

    else:
        raise ValueError(f"Unsupported file type: {extension}")
    
def process_page(idx, ocr_response=None):
    try:
        if ocr_response and hasattr(ocr_response, 'pages') and idx < len(ocr_response.pages):
            return ocr_response.pages[idx].markdown
        else:
            return f"Error: Page {idx + 1} not available in OCR response"
    except Exception as e:
        return f"Error processing page {idx + 1}: {e}"

def extract_text_from_pdf(pdf_bytes: bytes, advanced: bool = True):
    extracted_text = []

    if not advanced:
        # Simple text extraction using PyMuPDF (no OCR)
        try:
            with pymupdf.open(stream=pdf_bytes, filetype="pdf") as doc:
                for page in doc:
                    text = page.get_text()
                    extracted_text.append(text)
            return extracted_text
        except Exception as e:
            return [f"Error during simple text extraction: {e}"]

    # Advanced mode: Use Mistral OCR
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            total_pages = len(pdf.pages)
    except Exception:
        try:
            with pymupdf.open(stream=pdf_bytes, filetype="pdf") as doc:
                total_pages = len(doc)
        except Exception as e:
            return [f"Error getting total pages: {e}"]

    encoded_pdf = encode_pdf(pdf_bytes)

    try:
        response = client.ocr.process(
            model="mistral-ocr-latest",
            document={
                "type": "document_url",
                "document_url": f"data:application/pdf;base64,{encoded_pdf}"
            },
            include_image_base64=True
        )
    except Exception as e:
        return [f"Error during OCR processing: {e}"]

    for idx in range(total_pages):
        page_text = process_page(idx, ocr_response=response)
        extracted_text.append(page_text)
        
    return extracted_text


def create_chunks(directory_path: str):
    file_paths = [
        os.path.abspath(os.path.join(directory_path, f))
        for f in os.listdir(directory_path)
        if os.path.isfile(os.path.join(directory_path, f))
    ]
    
    Chunks = []
    
    for idx, file_path in enumerate(file_paths):
        filename = os.path.basename(file_path)
        extension = filename.lower().split('.')[-1]
        
        with open(file_path, "rb") as f:
            file_bytes = f.read()
        
        converted_pdf_bytes = convert_to_pdf(file_bytes, filename)
        print(f"Processing file: {filename}")
        
        # Decide mode based on file type
        if extension in {"txt", "md"}:
            pages = extract_text_from_pdf(converted_pdf_bytes, advanced=False)
        else:
            pages = extract_text_from_pdf(converted_pdf_bytes, advanced=True)
        
        for page_number, page in enumerate(pages, start=1):
            Chunks.append({
                "filename": filename,
                "page_number": page_number,
                "page_content": page
            })
    
    return Chunks
